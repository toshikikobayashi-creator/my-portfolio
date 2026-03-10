"""
ヒューマンアド 飲食店（ラーメン屋）向け提案書 PowerPoint生成スクリプト
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# === カラーパレット ===
COLOR_BG_DARK   = RGBColor(0x0F, 0x20, 0x27)   # 濃紺（背景メイン）
COLOR_BG_MID    = RGBColor(0x1A, 0x3A, 0x4A)   # 中間紺
COLOR_ACCENT    = RGBColor(0xE8, 0x6B, 0x1A)   # オレンジ（ラーメン・アクセント）
COLOR_ACCENT2   = RGBColor(0xFF, 0xC8, 0x5C)   # ゴールド（強調）
COLOR_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_LIGHT     = RGBColor(0xE8, 0xF4, 0xFF)   # 薄青（テキスト）
COLOR_GRAY      = RGBColor(0x8A, 0x9B, 0xB0)    # グレー
COLOR_SUCCESS   = RGBColor(0x27, 0xAE, 0x60)   # グリーン（ROI等）
COLOR_TABLE_HDR = RGBColor(0xE8, 0x6B, 0x1A)
COLOR_TABLE_ROW = RGBColor(0x12, 0x2A, 0x38)
COLOR_TABLE_ALT = RGBColor(0x1A, 0x3A, 0x4A)


def set_bg(slide, color: RGBColor):
    """スライド背景色を設定"""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(0)):
    """矩形シェイプを追加"""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             font_size=18, bold=False, color=COLOR_WHITE,
             align=PP_ALIGN.LEFT, italic=False):
    """テキストボックスを追加"""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Hiragino Kaku Gothic ProN"
    return txBox


def add_multiline_text(slide, lines, left, top, width, height,
                        font_size=14, bold=False, color=COLOR_WHITE,
                        align=PP_ALIGN.LEFT, line_spacing=None):
    """複数行テキストボックスを追加（リスト形式）"""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Hiragino Kaku Gothic ProN"
    return txBox


def add_accent_bar(slide, left, top, width=0.06, height=0.35, color=COLOR_ACCENT):
    """アクセントバー（見出し左のライン装飾）"""
    add_rect(slide, left, top, width, height, fill_color=color)


# ============================================================
# スライド生成関数
# ============================================================

def slide_01_cover(prs):
    """スライド1: 表紙"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白
    set_bg(slide, COLOR_BG_DARK)

    # 上部グラデーション風装飾帯
    add_rect(slide, 0, 0, 13.33, 0.08, fill_color=COLOR_ACCENT)

    # 左側縦ライン装飾
    add_rect(slide, 0.5, 0.8, 0.06, 5.0, fill_color=COLOR_ACCENT)

    # サービスロゴ風テキスト
    add_text(slide, "HUMAN AD", 0.7, 0.85, 5, 0.8,
             font_size=11, bold=True, color=COLOR_ACCENT)

    # メインタイトル
    add_text(slide, "集客力を\n歩かせよう。", 0.7, 1.5, 8, 2.2,
             font_size=44, bold=True, color=COLOR_WHITE)

    # サブタイトル
    add_text(slide, "ヒューマンアド 広告配信 ご提案書", 0.7, 3.6, 9, 0.6,
             font_size=20, bold=False, color=COLOR_ACCENT2)

    # 区切り線
    add_rect(slide, 0.7, 4.25, 5.5, 0.04, fill_color=COLOR_ACCENT)

    # クライアント・日付情報
    add_multiline_text(slide, [
        "ご提案先：ラーメン店 御中",
        "提案日　：2026年3月10日",
        "有効期限：2026年4月9日（30日間）",
        "提案元　：Backstage Inc. ヒューマンアド事業部",
    ], 0.7, 4.4, 8, 1.8, font_size=13, color=COLOR_LIGHT)

    # 右側イメージテキスト（装飾）
    add_rect(slide, 9.5, 1.2, 3.3, 4.5, fill_color=COLOR_BG_MID,
             line_color=COLOR_ACCENT, line_width=Pt(1.5))
    add_text(slide, "📡", 10.1, 1.6, 2, 1.2, font_size=52, align=PP_ALIGN.CENTER)
    add_text(slide, "32インチ LED\nバックパック広告", 9.6, 2.9, 3.2, 0.9,
             font_size=13, bold=True, color=COLOR_ACCENT2, align=PP_ALIGN.CENTER)
    add_multiline_text(slide, [
        "✓ 渋谷・新宿等 繁華街を歩行",
        "✓ 最大30,000人/時間へリーチ",
        "✓ アドトラックの1/10の低価格",
    ], 9.6, 3.85, 3.2, 1.6, font_size=11, color=COLOR_LIGHT)

    # 下部スライド番号エリア
    add_rect(slide, 0, 7.1, 13.33, 0.4, fill_color=COLOR_BG_MID)
    add_text(slide, "Backstage Inc.  |  Human Ad  |  Confidential", 0.3, 7.12, 8, 0.3,
             font_size=9, color=COLOR_GRAY, align=PP_ALIGN.LEFT)
    add_text(slide, "1", 12.8, 7.12, 0.4, 0.3, font_size=9, color=COLOR_GRAY, align=PP_ALIGN.RIGHT)


def slide_02_agenda(prs):
    """スライド2: 目次"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, COLOR_BG_DARK)
    add_rect(slide, 0, 0, 13.33, 0.08, fill_color=COLOR_ACCENT)

    add_accent_bar(slide, 0.5, 0.35)
    add_text(slide, "本日のご提案内容", 0.72, 0.3, 10, 0.55,
             font_size=22, bold=True, color=COLOR_WHITE)

    items = [
        ("01", "貴店の課題と機会",         "集客コストと効果の課題を整理"),
        ("02", "ヒューマンアドとは",       "サービス概要・強み・仕様"),
        ("03", "他媒体との比較",           "アドトラック・チラシとの違い"),
        ("04", "ご提案プランと料金",       "飲食店向けおすすめプラン"),
        ("05", "ROIシミュレーション",     "客単価¥1,000での収益予測"),
        ("06", "実施スケジュール",         "最短翌日から配信開始"),
        ("07", "次のステップ",             "お申し込みまでの流れ"),
    ]

    for i, (num, title, desc) in enumerate(items):
        y = 1.1 + i * 0.82
        # 番号ボックス
        add_rect(slide, 0.5, y, 0.55, 0.55, fill_color=COLOR_ACCENT)
        add_text(slide, num, 0.5, y + 0.02, 0.55, 0.5,
                 font_size=14, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
        # タイトル
        add_text(slide, title, 1.18, y + 0.02, 4.5, 0.35,
                 font_size=15, bold=True, color=COLOR_WHITE)
        # 説明
        add_text(slide, desc, 1.18, y + 0.3, 6.0, 0.3,
                 font_size=11, color=COLOR_LIGHT)
        # 区切り線
        if i < len(items) - 1:
            add_rect(slide, 0.5, y + 0.68, 12.3, 0.01, fill_color=COLOR_BG_MID)

    add_rect(slide, 0, 7.1, 13.33, 0.4, fill_color=COLOR_BG_MID)
    add_text(slide, "Backstage Inc.  |  Human Ad  |  Confidential", 0.3, 7.12, 8, 0.3,
             font_size=9, color=COLOR_GRAY)
    add_text(slide, "2", 12.8, 7.12, 0.4, 0.3, font_size=9, color=COLOR_GRAY, align=PP_ALIGN.RIGHT)


def slide_03_problem(prs):
    """スライド3: 課題整理"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, COLOR_BG_DARK)
    add_rect(slide, 0, 0, 13.33, 0.08, fill_color=COLOR_ACCENT)

    add_accent_bar(slide, 0.5, 0.35)
    add_text(slide, "01  貴店の課題と機会", 0.72, 0.3, 10, 0.55,
             font_size=22, bold=True, color=COLOR_WHITE)

    # 左：課題
    add_rect(slide, 0.4, 1.05, 5.9, 5.7, fill_color=COLOR_BG_MID,
             line_color=RGBColor(0xE0, 0x40, 0x40), line_width=Pt(1.5))
    add_text(slide, "❌  現状の課題", 0.55, 1.12, 5.6, 0.45,
             font_size=15, bold=True, color=RGBColor(0xFF, 0x6B, 0x6B))
    problems = [
        "チラシは配っても効果が見えにくい",
        "SNS広告は設定が複雑でコストが読めない",
        "ランチタイムに近隣からの来店が少ない",
        "「知っているけど行ったことがない」\n層への接点がない",
        "繁忙期だけ集中して告知したいが\n短期対応できる媒体がない",
    ]
    for i, prob in enumerate(problems):
        add_text(slide, f"▶  {prob}", 0.55, 1.65 + i * 0.9, 5.6, 0.8,
                 font_size=12, color=COLOR_LIGHT)

    # 右：機会
    add_rect(slide, 6.9, 1.05, 5.9, 5.7, fill_color=COLOR_BG_MID,
             line_color=COLOR_SUCCESS, line_width=Pt(1.5))
    add_text(slide, "✅  ヒューマンアドで解決", 7.05, 1.12, 5.6, 0.45,
             font_size=15, bold=True, color=COLOR_SUCCESS)
    solutions = [
        "通行人に「直接」リーチ。視覚的\nインパクトで記憶に残る",
        "最短翌日から配信。効果が出なければ\n即停止できる柔軟さ",
        "ランチ・夕方など時間帯を指定して\n近隣の人に集中アプローチ",
        "SNS投稿を誘発する「撮られる広告」\nで二次拡散も期待できる",
        "1日¥15,000〜の低価格で試せる\nチラシ代わりの即効性媒体",
    ]
    for i, sol in enumerate(solutions):
        add_text(slide, f"▶  {sol}", 7.05, 1.65 + i * 0.9, 5.6, 0.8,
                 font_size=12, color=COLOR_LIGHT)

    add_rect(slide, 0, 7.1, 13.33, 0.4, fill_color=COLOR_BG_MID)
    add_text(slide, "Backstage Inc.  |  Human Ad  |  Confidential", 0.3, 7.12, 8, 0.3,
             font_size=9, color=COLOR_GRAY)
    add_text(slide, "3", 12.8, 7.12, 0.4, 0.3, font_size=9, color=COLOR_GRAY, align=PP_ALIGN.RIGHT)


def slide_04_service(prs):
    """スライド4: サービス概要"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, COLOR_BG_DARK)
    add_rect(slide, 0, 0, 13.33, 0.08, fill_color=COLOR_ACCENT)

    add_accent_bar(slide, 0.5, 0.35)
    add_text(slide, "02  ヒューマンアドとは", 0.72, 0.3, 10, 0.55,
             font_size=22, bold=True, color=COLOR_WHITE)

    # 中央：サービス説明ボックス
    add_rect(slide, 0.4, 1.05, 12.5, 1.3, fill_color=COLOR_BG_MID)
    add_text(slide,
             "スタッフが32インチの大型LEDディスプレイを背負って、ターゲットエリアを歩行する\n"
             "まったく新しい「モバイル × デジタル × 人力」の屋外広告サービスです。",
             0.6, 1.1, 12.1, 1.2, font_size=14, color=COLOR_WHITE)

    # 3つの強みカード
    cards = [
        ("📍", "モバイル", "車が入れない\n歩行者エリアも配信可能\nモール・公園・祭会場でも対応"),
        ("⚡", "即時対応", "最短翌日から配信開始\n思い立ったらすぐ実施\nイベント直前でも間に合う"),
        ("💰", "低価格", "1日¥15,000〜\nアドトラックの1/10以下\nスモールスタートで検証可能"),
    ]
    for i, (icon, title, desc) in enumerate(cards):
        x = 0.4 + i * 4.3
        add_rect(slide, x, 2.55, 4.0, 3.2, fill_color=COLOR_BG_MID,
                 line_color=COLOR_ACCENT, line_width=Pt(1.5))
        add_text(slide, icon, x, 2.7, 4.0, 0.8,
                 font_size=32, align=PP_ALIGN.CENTER)
        add_text(slide, title, x, 3.5, 4.0, 0.5,
                 font_size=17, bold=True, color=COLOR_ACCENT2, align=PP_ALIGN.CENTER)
        add_text(slide, desc, x + 0.2, 4.05, 3.6, 1.6,
                 font_size=12, color=COLOR_LIGHT, align=PP_ALIGN.CENTER)

    # 仕様情報
    add_rect(slide, 0.4, 5.95, 12.5, 0.95, fill_color=RGBColor(0x0A, 0x18, 0x22))
    add_text(slide, "【機材仕様】  Jiatean JTA-WT32-BA  |  32インチ LEDディスプレイ  |  高輝度・全天候対応  |  動画・静止画スライドショー対応",
             0.6, 6.05, 12.1, 0.7, font_size=11, color=COLOR_ACCENT2)

    add_rect(slide, 0, 7.1, 13.33, 0.4, fill_color=COLOR_BG_MID)
    add_text(slide, "Backstage Inc.  |  Human Ad  |  Confidential", 0.3, 7.12, 8, 0.3,
             font_size=9, color=COLOR_GRAY)
    add_text(slide, "4", 12.8, 7.12, 0.4, 0.3, font_size=9, color=COLOR_GRAY, align=PP_ALIGN.RIGHT)


def slide_05_comparison(prs):
    """スライド5: 他媒体との比較"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, COLOR_BG_DARK)
    add_rect(slide, 0, 0, 13.33, 0.08, fill_color=COLOR_ACCENT)

    add_accent_bar(slide, 0.5, 0.35)
    add_text(slide, "03  他媒体との比較", 0.72, 0.3, 10, 0.55,
             font_size=22, bold=True, color=COLOR_WHITE)

    # 比較表ヘッダー
    headers = ["比較項目", "ヒューマンアド", "アドトラック", "チラシ配布", "SNS広告"]
    col_x = [0.35, 2.65, 5.15, 7.55, 9.95]
    col_w = [2.2, 2.4, 2.3, 2.3, 2.9]

    for i, (hdr, x, w) in enumerate(zip(headers, col_x, col_w)):
        bg = COLOR_ACCENT if i == 1 else COLOR_TABLE_HDR if i == 0 else COLOR_BG_MID
        add_rect(slide, x, 1.05, w, 0.5, fill_color=bg)
        add_text(slide, hdr, x, 1.07, w, 0.45,
                 font_size=12, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)

    rows = [
        ("1日あたり費用",    "¥15,000〜",      "¥150,000〜",   "¥5,000〜",      "¥3,000〜"),
        ("展開スピード",     "翌日〜即日",      "1週間〜",       "2〜3日",         "即日"),
        ("ターゲット精度",   "エリア+時間帯",   "エリア指定",    "近隣限定",       "属性指定"),
        ("視覚的インパクト", "◎ 大（32インチ）","◎ 大（車体）",  "△ 小",          "△ 小"),
        ("SNS拡散効果",     "◎ 高",           "△ 中",         "✕ なし",         "◎ 高"),
        ("規制リスク",       "◎ 対象外",       "✕ 強化中",      "○ 低",          "○ 低"),
        ("歩行者エリア対応", "◎ 可能",         "✕ 不可",        "○ 可能",         "— "),
        ("効果の可視性",     "○ 推定リーチ",   "○ 推定リーチ",  "✕ 不明",        "◎ 数値管理"),
    ]

    for r, (row_data) in enumerate(rows):
        y = 1.6 + r * 0.66
        bg = COLOR_TABLE_ROW if r % 2 == 0 else COLOR_TABLE_ALT
        for i, (cell, x, w) in enumerate(zip(row_data, col_x, col_w)):
            cell_bg = RGBColor(0x1E, 0x42, 0x1E) if i == 1 else bg  # ヒューマンアド列を強調
            add_rect(slide, x, y, w, 0.62, fill_color=cell_bg)
            color = COLOR_ACCENT2 if i == 1 else COLOR_WHITE
            bold = i == 1
            add_text(slide, cell, x + 0.05, y + 0.06, w - 0.1, 0.5,
                     font_size=11, bold=bold, color=color, align=PP_ALIGN.CENTER)

    add_rect(slide, 0, 7.1, 13.33, 0.4, fill_color=COLOR_BG_MID)
    add_text(slide, "Backstage Inc.  |  Human Ad  |  Confidential", 0.3, 7.12, 8, 0.3,
             font_size=9, color=COLOR_GRAY)
    add_text(slide, "5", 12.8, 7.12, 0.4, 0.3, font_size=9, color=COLOR_GRAY, align=PP_ALIGN.RIGHT)


def slide_06_plan(prs):
    """スライド6: 提案プランと料金"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, COLOR_BG_DARK)
    add_rect(slide, 0, 0, 13.33, 0.08, fill_color=COLOR_ACCENT)

    add_accent_bar(slide, 0.5, 0.35)
    add_text(slide, "04  ご提案プランと料金", 0.72, 0.3, 10, 0.55,
             font_size=22, bold=True, color=COLOR_WHITE)

    # 3プランカード
    plans = [
        {
            "name": "お試しプラン",
            "tag": "まず試したい方に",
            "price": "¥15,000",
            "unit": "（税抜 / 1日）",
            "details": [
                "スタッフ 1名",
                "配信時間 2時間",
                "エリア：店舗周辺500m圏内",
                "コンテンツ：お客様ご提供",
                "推定リーチ：約4,000人",
            ],
            "highlight": False,
        },
        {
            "name": "ランチ集客プラン",
            "tag": "★ 最もおすすめ",
            "price": "¥25,000",
            "unit": "（税抜 / 1日）",
            "details": [
                "スタッフ 1名",
                "配信時間 4時間（11時〜15時）",
                "エリア：駅周辺〜店舗導線",
                "コンテンツ：弊社制作サポート付き",
                "推定リーチ：約10,000人",
            ],
            "highlight": True,
        },
        {
            "name": "週末フルプラン",
            "tag": "週末集中で最大効果",
            "price": "¥45,000",
            "unit": "（税抜 / 1日）",
            "details": [
                "スタッフ 2名",
                "配信時間 4時間",
                "エリア：複数ルートで同時配信",
                "コンテンツ：弊社制作",
                "推定リーチ：約20,000人",
            ],
            "highlight": False,
        },
    ]

    for i, plan in enumerate(plans):
        x = 0.35 + i * 4.3
        border_color = COLOR_ACCENT2 if plan["highlight"] else COLOR_ACCENT
        border_w = Pt(2.5) if plan["highlight"] else Pt(1.0)
        add_rect(slide, x, 1.05, 4.0, 5.8,
                 fill_color=RGBColor(0x1E, 0x42, 0x1E) if plan["highlight"] else COLOR_BG_MID,
                 line_color=border_color, line_width=border_w)

        # おすすめバッジ
        if plan["highlight"]:
            add_rect(slide, x + 0.8, 1.0, 2.4, 0.38, fill_color=COLOR_ACCENT2)
            add_text(slide, plan["tag"], x + 0.8, 1.01, 2.4, 0.35,
                     font_size=11, bold=True, color=COLOR_BG_DARK, align=PP_ALIGN.CENTER)
        else:
            add_text(slide, plan["tag"], x, 1.12, 4.0, 0.3,
                     font_size=10, color=COLOR_LIGHT, align=PP_ALIGN.CENTER)

        add_text(slide, plan["name"], x, 1.5, 4.0, 0.5,
                 font_size=16, bold=True,
                 color=COLOR_ACCENT2 if plan["highlight"] else COLOR_WHITE,
                 align=PP_ALIGN.CENTER)
        add_text(slide, plan["price"], x, 2.05, 4.0, 0.75,
                 font_size=32, bold=True, color=COLOR_ACCENT2, align=PP_ALIGN.CENTER)
        add_text(slide, plan["unit"], x, 2.75, 4.0, 0.3,
                 font_size=10, color=COLOR_LIGHT, align=PP_ALIGN.CENTER)

        # 区切り線
        add_rect(slide, x + 0.2, 3.1, 3.6, 0.04, fill_color=COLOR_ACCENT)

        for j, detail in enumerate(plan["details"]):
            add_text(slide, f"✓  {detail}", x + 0.2, 3.2 + j * 0.5, 3.6, 0.45,
                     font_size=11, color=COLOR_LIGHT)

    # 注記
    add_text(slide, "※ 表示料金はすべて税抜。消費税10%が別途かかります。コンテンツ制作費は別途お見積り。",
             0.35, 7.0, 12.5, 0.3, font_size=9, color=COLOR_GRAY)

    add_rect(slide, 0, 7.1, 13.33, 0.4, fill_color=COLOR_BG_MID)
    add_text(slide, "Backstage Inc.  |  Human Ad  |  Confidential", 0.3, 7.12, 8, 0.3,
             font_size=9, color=COLOR_GRAY)
    add_text(slide, "6", 12.8, 7.12, 0.4, 0.3, font_size=9, color=COLOR_GRAY, align=PP_ALIGN.RIGHT)


def slide_07_roi(prs):
    """スライド7: ROIシミュレーション"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, COLOR_BG_DARK)
    add_rect(slide, 0, 0, 13.33, 0.08, fill_color=COLOR_ACCENT)

    add_accent_bar(slide, 0.5, 0.35)
    add_text(slide, "05  ROIシミュレーション（客単価 ¥1,000 ラーメン店想定）", 0.72, 0.3, 12, 0.55,
             font_size=20, bold=True, color=COLOR_WHITE)

    # 計算フロー（左側）
    add_text(slide, "■ ランチ集客プランの場合（¥25,000/日）", 0.4, 1.05, 8, 0.4,
             font_size=13, bold=True, color=COLOR_ACCENT2)

    steps = [
        ("推定通行人数", "駅周辺エリア（平日ランチ）", "10,000人 / 4時間"),
        ("視認率",       "LEDディスプレイへの目視率（推定）", "× 20%  →  2,000人"),
        ("有効記憶率",   "広告内容を記憶した割合（推定）",   "× 25%  →  500人"),
        ("来店転換率",   "記憶後に来店した割合（飲食業推定）","× 3%   →  15人/日"),
        ("売上効果",     "来店者 × 客単価 ¥1,000",           "= ¥15,000 / 日"),
    ]

    arrows = ["↓", "↓", "↓", "↓"]
    for i, (title, desc, value) in enumerate(steps):
        y = 1.55 + i * 0.95
        add_rect(slide, 0.4, y, 8.3, 0.75, fill_color=COLOR_BG_MID)
        add_text(slide, title, 0.55, y + 0.04, 2.8, 0.35,
                 font_size=12, bold=True, color=COLOR_ACCENT2)
        add_text(slide, desc, 0.55, y + 0.36, 4.5, 0.3,
                 font_size=10, color=COLOR_LIGHT)
        add_text(slide, value, 5.8, y + 0.12, 2.8, 0.5,
                 font_size=14, bold=True, color=COLOR_WHITE, align=PP_ALIGN.RIGHT)
        if i < len(arrows):
            add_text(slide, arrows[i], 4.2, y + 0.75, 0.5, 0.3,
                     font_size=12, color=COLOR_ACCENT, align=PP_ALIGN.CENTER)

    # ROI結果ボックス（右側）
    add_rect(slide, 9.0, 1.05, 3.95, 5.75, fill_color=RGBColor(0x1E, 0x42, 0x1E),
             line_color=COLOR_SUCCESS, line_width=Pt(2))
    add_text(slide, "ROI 試算", 9.0, 1.12, 3.95, 0.45,
             font_size=14, bold=True, color=COLOR_SUCCESS, align=PP_ALIGN.CENTER)

    roi_items = [
        ("広告費用", "¥25,000"),
        ("1日売上効果", "¥15,000"),
        ("3日間実施", "¥45,000"),
        ("3日間効果", "¥45,000"),
        ("ROI（3日）", "80%"),
        ("損益分岐", "来店 25人〜"),
    ]
    for i, (label, val) in enumerate(roi_items):
        y = 1.65 + i * 0.8
        add_text(slide, label, 9.1, y, 2.2, 0.4, font_size=11, color=COLOR_LIGHT)
        color = COLOR_ACCENT2 if "ROI" in label or "損益" in label else COLOR_WHITE
        add_text(slide, val, 11.0, y, 1.8, 0.4,
                 font_size=13, bold=True, color=color, align=PP_ALIGN.RIGHT)
        add_rect(slide, 9.1, y + 0.44, 3.7, 0.02, fill_color=COLOR_BG_MID)

    # 注記
    add_text(slide, "※ 数値は推定値。エリア・時間帯・コンテンツ品質により変動します。",
             0.4, 7.0, 12.5, 0.3, font_size=9, color=COLOR_GRAY)

    add_rect(slide, 0, 7.1, 13.33, 0.4, fill_color=COLOR_BG_MID)
    add_text(slide, "Backstage Inc.  |  Human Ad  |  Confidential", 0.3, 7.12, 8, 0.3,
             font_size=9, color=COLOR_GRAY)
    add_text(slide, "7", 12.8, 7.12, 0.4, 0.3, font_size=9, color=COLOR_GRAY, align=PP_ALIGN.RIGHT)


def slide_08_schedule(prs):
    """スライド8: 実施スケジュール"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, COLOR_BG_DARK)
    add_rect(slide, 0, 0, 13.33, 0.08, fill_color=COLOR_ACCENT)

    add_accent_bar(slide, 0.5, 0.35)
    add_text(slide, "06  実施スケジュール", 0.72, 0.3, 10, 0.55,
             font_size=22, bold=True, color=COLOR_WHITE)

    add_text(slide, "最短翌日から配信開始できます", 0.5, 1.0, 12.0, 0.5,
             font_size=16, bold=True, color=COLOR_ACCENT2)

    steps = [
        ("D-3以上前", "お申し込み・ヒアリング",
         "・希望エリア・日時・配信内容をご確認\n・コンテンツ素材のご提供またはご相談"),
        ("D-2",      "コンテンツ確認・最終調整",
         "・配信データの確認・テスト表示\n・ルート・タイムラインの最終確定"),
        ("D-1",      "機材・スタッフ準備",
         "・機材チェック・充電\n・スタッフへのブリーフィング"),
        ("当日(D)",   "配信実施",
         "・指定エリアを歩行しながら配信\n・状況報告（写真・動画をSlack/メール共有）"),
        ("D+1",      "効果測定レポート提出",
         "・推定リーチ数・配信実績レポートをご送付\n・次回施策のご提案"),
    ]

    for i, (day, title, desc) in enumerate(steps):
        y = 1.6 + i * 1.0
        # タイムラインライン
        add_rect(slide, 1.5, y + 0.2, 0.04, 1.0,
                 fill_color=COLOR_ACCENT if i < len(steps) - 1 else RGBColor(0x00,0x00,0x00))
        # 丸
        add_rect(slide, 1.35, y + 0.08, 0.35, 0.35, fill_color=COLOR_ACCENT)
        # 日付
        add_text(slide, day, 0.4, y + 0.1, 0.9, 0.35,
                 font_size=10, bold=True, color=COLOR_ACCENT2, align=PP_ALIGN.RIGHT)
        # タイトル
        add_text(slide, title, 1.85, y + 0.06, 4.0, 0.35,
                 font_size=14, bold=True, color=COLOR_WHITE)
        # 説明
        add_text(slide, desc, 1.85, y + 0.42, 10.5, 0.5,
                 font_size=11, color=COLOR_LIGHT)

    add_rect(slide, 0, 7.1, 13.33, 0.4, fill_color=COLOR_BG_MID)
    add_text(slide, "Backstage Inc.  |  Human Ad  |  Confidential", 0.3, 7.12, 8, 0.3,
             font_size=9, color=COLOR_GRAY)
    add_text(slide, "8", 12.8, 7.12, 0.4, 0.3, font_size=9, color=COLOR_GRAY, align=PP_ALIGN.RIGHT)


def slide_09_nextstep(prs):
    """スライド9: 次のステップ"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, COLOR_BG_DARK)
    add_rect(slide, 0, 0, 13.33, 0.08, fill_color=COLOR_ACCENT)

    add_accent_bar(slide, 0.5, 0.35)
    add_text(slide, "07  次のステップ", 0.72, 0.3, 10, 0.55,
             font_size=22, bold=True, color=COLOR_WHITE)

    # CTAボックス
    add_rect(slide, 0.5, 1.05, 12.3, 1.2, fill_color=RGBColor(0x1E, 0x42, 0x1E),
             line_color=COLOR_SUCCESS, line_width=Pt(1.5))
    add_text(slide, "今すぐ「お試しプラン（¥15,000〜）」からスタートできます",
             0.7, 1.2, 11.9, 0.5, font_size=17, bold=True, color=COLOR_ACCENT2, align=PP_ALIGN.CENTER)
    add_text(slide, "コンテンツは貴店のInstagram画像・メニュー写真でOK！初回は弊社がサポートします。",
             0.7, 1.65, 11.9, 0.4, font_size=12, color=COLOR_LIGHT, align=PP_ALIGN.CENTER)

    # ステップ
    steps = [
        ("STEP 1", "ご連絡", "本資料をご確認の上、下記連絡先へご連絡ください。"),
        ("STEP 2", "ヒアリング", "ご希望日時・エリア・コンテンツについてご確認します（30分程度）。"),
        ("STEP 3", "お見積り", "正式なお見積書を発行します（無料）。"),
        ("STEP 4", "配信開始", "ご入金確認後、最短翌日から配信をスタートします。"),
    ]

    for i, (step, title, desc) in enumerate(steps):
        x = 0.5 + i * 3.2
        add_rect(slide, x, 2.5, 3.0, 2.8, fill_color=COLOR_BG_MID,
                 line_color=COLOR_ACCENT, line_width=Pt(1))
        add_rect(slide, x, 2.5, 3.0, 0.45, fill_color=COLOR_ACCENT)
        add_text(slide, step, x, 2.54, 3.0, 0.38,
                 font_size=13, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, title, x, 3.0, 3.0, 0.45,
                 font_size=15, bold=True, color=COLOR_ACCENT2, align=PP_ALIGN.CENTER)
        add_text(slide, desc, x + 0.15, 3.5, 2.7, 1.6,
                 font_size=11, color=COLOR_LIGHT)

    # 連絡先ボックス
    add_rect(slide, 0.5, 5.55, 12.3, 1.3, fill_color=COLOR_BG_MID,
             line_color=COLOR_ACCENT, line_width=Pt(1))
    add_text(slide, "お問い合わせ先", 0.7, 5.62, 4, 0.35,
             font_size=13, bold=True, color=COLOR_ACCENT2)
    add_multiline_text(slide, [
        "Backstage Inc.  ヒューマンアド事業部",
        "Email: info@backstage-inc.example.com   |   Tel: 03-XXXX-XXXX",
        "営業時間: 平日 10:00〜18:00",
    ], 0.7, 5.98, 11.8, 0.85, font_size=12, color=COLOR_LIGHT)

    # 下部フッター
    add_rect(slide, 0, 7.1, 13.33, 0.4, fill_color=COLOR_BG_MID)
    add_text(slide, "Backstage Inc.  |  Human Ad  |  Confidential", 0.3, 7.12, 8, 0.3,
             font_size=9, color=COLOR_GRAY)
    add_text(slide, "9", 12.8, 7.12, 0.4, 0.3, font_size=9, color=COLOR_GRAY, align=PP_ALIGN.RIGHT)


def slide_10_appendix(prs):
    """スライド10: 料金オプション一覧（付録）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, COLOR_BG_DARK)
    add_rect(slide, 0, 0, 13.33, 0.08, fill_color=COLOR_ACCENT)

    add_accent_bar(slide, 0.5, 0.35)
    add_text(slide, "Appendix  料金オプション一覧", 0.72, 0.3, 10, 0.55,
             font_size=22, bold=True, color=COLOR_WHITE)

    # オプションテーブル
    headers = ["オプション", "料金（税抜）", "備考"]
    col_x = [0.4, 6.5, 9.5]
    col_w = [6.0, 2.8, 3.3]

    for i, (h, x, w) in enumerate(zip(headers, col_x, col_w)):
        add_rect(slide, x, 1.05, w, 0.45, fill_color=COLOR_TABLE_HDR)
        add_text(slide, h, x, 1.1, w, 0.38,
                 font_size=13, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)

    opt_rows = [
        ("エリア指定（渋谷・新宿・原宿・銀座等）", "¥5,000/日", "都心主要エリアのみ"),
        ("動画コンテンツ制作",                     "¥30,000〜", "15〜30秒ループ動画"),
        ("静止画スライドショー制作",               "¥15,000〜", "メニュー写真+テキスト"),
        ("SNS投稿レポート",                        "¥10,000",   "1日あたり拡散状況レポート"),
        ("効果測定レポート",                       "¥20,000",   "推定リーチ数・接触者数レポート"),
        ("急ぎ対応（翌日手配）",                   "¥8,000",    "通常3営業日前→翌日"),
        ("キャンセル料（2〜3日前）",               "見積の50%", "悪天候は免除"),
        ("キャンセル料（前日）",                   "見積の80%", "悪天候は免除"),
    ]

    for r, (item, price, note) in enumerate(opt_rows):
        y = 1.55 + r * 0.63
        bg = COLOR_TABLE_ROW if r % 2 == 0 else COLOR_TABLE_ALT
        for val, x, w in zip([item, price, note], col_x, col_w):
            add_rect(slide, x, y, w, 0.6, fill_color=bg)
            add_text(slide, val, x + 0.1, y + 0.1, w - 0.2, 0.42,
                     font_size=11, color=COLOR_WHITE)

    add_rect(slide, 0, 7.1, 13.33, 0.4, fill_color=COLOR_BG_MID)
    add_text(slide, "Backstage Inc.  |  Human Ad  |  Confidential", 0.3, 7.12, 8, 0.3,
             font_size=9, color=COLOR_GRAY)
    add_text(slide, "10", 12.8, 7.12, 0.4, 0.3, font_size=9, color=COLOR_GRAY, align=PP_ALIGN.RIGHT)


# ============================================================
# メイン処理
# ============================================================

def main():
    prs = Presentation()
    # スライドサイズをワイド（16:9）に設定
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    slide_01_cover(prs)
    slide_02_agenda(prs)
    slide_03_problem(prs)
    slide_04_service(prs)
    slide_05_comparison(prs)
    slide_06_plan(prs)
    slide_07_roi(prs)
    slide_08_schedule(prs)
    slide_09_nextstep(prs)
    slide_10_appendix(prs)

    output_path = "/Users/macuser/my-portfolio/sales/proposals/proposal_ramen_humanad_20260310.pptx"
    prs.save(output_path)
    print(f"✅ 保存完了: {output_path}")
    print(f"   スライド数: {len(prs.slides)}枚")


if __name__ == "__main__":
    main()
